import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

def create_features_presentation():
    prs = Presentation()
    
    script_dir = os.path.dirname(os.path.abspath(__file__))
    ASSETS_DIR = os.path.join(script_dir, "assets")
    
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # blank slide layout

    # Color Palette
    NAVY = RGBColor(15, 76, 129)          # #0F4C81 - Primary Deep
    CYAN = RGBColor(0, 180, 216)          # #00B4D8 - Accent Teal/Cyan
    DARK_BG = RGBColor(11, 25, 44)        # #0B192C - Premium Dark BG
    LIGHT_BG = RGBColor(248, 250, 252)    # #F8FAFC - Clean Light BG
    CARD_BG = RGBColor(255, 255, 255)     # #FFFFFF
    TEXT_MAIN = RGBColor(30, 41, 59)      # #1E293B
    TEXT_MUTED = RGBColor(100, 116, 139)  # #64748B
    ACCENT_BLUE = RGBColor(224, 242, 254) # #E0F2FE

    BADGE_COLORS = [
        RGBColor(15, 76, 129),    # Deep Navy
        RGBColor(0, 150, 199),   # Vibrant Cyan/Teal
        RGBColor(16, 124, 65),   # Emerald Green
        RGBColor(217, 119, 6),   # Amber Gold
        RGBColor(124, 58, 237),  # Royal Purple
    ]

    def add_bg(slide, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_header(slide, title_text, category_text="TrainRouter BD", dark_theme=False):
        # Header bar badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.9), Inches(0.38))
        badge.fill.solid()
        badge.fill.fore_color.rgb = CYAN if dark_theme else NAVY
        badge.line.fill.background()
        tf = badge.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"►  {category_text.upper()}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = DARK_BG if dark_theme else RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Title
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.82), Inches(9.5), Inches(0.75))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255) if dark_theme else NAVY

        # Top Right Brand Badge
        brand_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.7), Inches(0.4), Inches(1.8), Inches(0.38))
        brand_badge.fill.solid()
        brand_badge.fill.fore_color.rgb = RGBColor(30, 41, 59) if dark_theme else RGBColor(241, 245, 249)
        brand_badge.line.color.rgb = CYAN if dark_theme else NAVY
        brand_badge.line.width = Pt(1.5)
        tf_b = brand_badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = "TrainRouter BD"
        p_b.font.size = Pt(11)
        p_b.font.bold = True
        p_b.font.color.rgb = CYAN if dark_theme else NAVY
        p_b.alignment = PP_ALIGN.CENTER

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    def add_feature_slide(slide, title_text, category_text, points_list, img_path):
        add_bg(slide, LIGHT_BG)
        add_header(slide, title_text, category_text)

        # Left Content Container Card
        add_card(slide, Inches(0.8), Inches(1.55), Inches(5.6), Inches(5.45), CARD_BG, border_color=NAVY)
        
        tb = slide.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(5.3), Inches(5.15))
        tf = tb.text_frame
        tf.word_wrap = True
        
        first = True
        for idx, item in enumerate(points_list):
            if len(item) == 3:
                icon_sym, t_text, d_text = item
            elif len(item) == 2:
                icon_sym = "►"
                t_text, d_text = item
            else:
                icon_sym = "►"
                t_text, d_text = f"Feature Point 0{idx+1}", item[0]

            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            
            p.text = f"{icon_sym}  {t_text}"
            p.font.size = Pt(15.5)
            p.font.bold = True
            p.font.color.rgb = BADGE_COLORS[idx % len(BADGE_COLORS)]
            p.space_before = Pt(6) if idx > 0 else Pt(0)
            
            p_desc = tf.add_paragraph()
            p_desc.text = d_text
            p_desc.font.size = Pt(13)
            p_desc.font.color.rgb = TEXT_MAIN
            p_desc.space_after = Pt(6)

        # Right Image Frame Card (Expanded Wide Frame)
        card_left = Inches(6.6)
        card_top = Inches(1.55)
        card_width = Inches(5.9)
        card_height = Inches(5.45)
        add_card(slide, card_left, card_top, card_width, card_height, RGBColor(241, 245, 249), border_color=CYAN)
        
        if img_path and os.path.exists(img_path):
            with Image.open(img_path) as img:
                orig_w, orig_h = img.size
            aspect = orig_w / float(orig_h)
            
            # Calculate optimal wide dimensions preserving natural aspect ratio
            max_img_w = Inches(5.6)
            max_img_h = Inches(5.1)
            
            calc_h = max_img_w / aspect
            if calc_h <= max_img_h:
                final_w = max_img_w
                final_h = calc_h
            else:
                final_h = max_img_h
                final_w = final_h * aspect
                
            img_left = card_left + (card_width - final_w) / 2.0
            img_top = card_top + (card_height - final_h) / 2.0
            
            slide.shapes.add_picture(img_path, img_left, img_top, final_w, final_h)

    mockup_img_v2 = os.path.join(ASSETS_DIR, "train_app_mockup_v2.jpg")
    mockup_img = mockup_img_v2 if os.path.exists(mockup_img_v2) else os.path.join(ASSETS_DIR, "train_app_mockup.png")

    # -------------------------------------------------------------------------
    # SLIDE 1: System Features Showcase Overview
    # -------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    feats_overview = [
        ("🔍", "From-To Station Train Search", "Select origin, destination & date to view all active running trains."),
        ("📊", "Available Seat Counts per Class", "Live inventory breakdown for Shovon, Snigdha, AC Berth, and AC Cabin."),
        ("🧮", "Seat-Wise Fare Calculator (BDT ৳)", "Instant cost lookup for single seat & multi-passenger journeys."),
        ("📍", "Intermediate Stoppages & Timetables", "Arrival/departure times at major junction stations."),
        ("📌", "Pinned Ticket Store & Departure Alarm", "Store offline ticket reference & set departure alert notifications.")
    ]
    add_feature_slide(slide1, "System Features Showcase Overview", "Features Overview", feats_overview, mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 2: Feature 01 - From-To & Date-Based Search
    # -------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    feats1 = [
        ("🔎", "Smart Autocomplete Station Lookup", "Select From (Dhaka Kamalapur) and To (Chattogram, Sylhet, Cox's Bazar, Rajshahi)."),
        ("📅", "Journey Date Picker & Off-Day Filter", "Pick any travel date. System automatically checks weekly off-days and filters non-operating trains."),
        ("🏷️", "Available Train Counter Badge", "Instantly shows top header summary: '4 Intercity Trains Available on Selected Date'.")
    ]
    add_feature_slide(slide2, "From-To & Date-Based Train Search", "Feature 01", feats1, mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 3: Feature 02 - Date-Wise Available Train List & Seat Count
    # -------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    feats2 = [
        ("🚆", "Date-Wise Active Train Schedule", "Displays exact running intercity trains on selected date (Suborno 701, Parabat 709, Cox's Bazar Express 813)."),
        ("💺", "Real-Time Seat Inventory per Class", "Live seat count breakdown (Shovon: 120 seats, Snigdha: 45 seats, AC Berth: 12 seats)."),
        ("⏰", "Train Timetable Details", "Exact departure time, destination arrival time, and estimated total journey duration.")
    ]
    add_feature_slide(slide3, "Date-Wise Active Train Schedule & Seats", "Feature 02", feats2, mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 4: Feature 03 - Seat Categories & Specifications
    # -------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    seats_info = [
        ("💺", "Shovon Chair (S_CHAIR)", "Standard non-AC padded recliner seat with 1.0x base tariff rate."),
        ("❄️", "Snigdha (AC Chair)", "Premium air-conditioned recliner seating with 1.8x base tariff rate."),
        ("🌙", "AC Berth (AC_BATH)", "Night sleeper berth compartment for long-distance overnight travel (2.5x tariff)."),
        ("👑", "AC Cabin (AC_CABIN)", "Private lockable compartment cabin suitable for family travel (2.8x tariff).")
    ]
    f3_img = os.path.join(ASSETS_DIR, "feature_03_seat_types.png")
    add_feature_slide(slide4, "Available Seat Types & Categories", "Feature 03", seats_info, f3_img if os.path.exists(f3_img) else mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 5: Feature 04 - Seat-Wise Fare Calculator
    # -------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    fare_info = [
        ("📏", "Dynamic Distance Fare Calculation", "Automatically looks up route distance (KM) between selected stations."),
        ("🏷️", "Class Multipliers & VAT Addition", "Applies exact seat multiplier (1.0x to 2.8x) + 15% Bangladesh Govt VAT."),
        ("👨‍👩‍👧‍👦", "Group / Family Trip Costing", "Calculates total cost for N passengers dynamically without requiring ticket booking.")
    ]
    fare_user_img = os.path.join(ASSETS_DIR, "fare_calc_user_mockup.jpg")
    add_feature_slide(slide5, "Seat-Wise Fare Calculator (BDT ৳)", "Feature 04", fare_info, fare_user_img if os.path.exists(fare_user_img) else mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 6: Feature 05 - Stoppage & Timetable Planner
    # -------------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    stop_info = [
        ("🚉", "Stoppage Junction Timetables", "Detailed arrival & departure times at intermediate stations (Airport, Joydebpur, Akhaura)."),
        ("🌐", "Cross-Zone Rail Connectivity", "Route mapping connecting Dhaka, Rajshahi, Khulna, Sylhet, and Cox's Bazar."),
        ("🛤️", "Gauge Compatibility Indicator", "Clearly demarcates Meter Gauge (MG), Broad Gauge (BG), and Dual Gauge (DG) lines.")
    ]
    f5_img = os.path.join(ASSETS_DIR, "feature_05_route_stoppages.png")
    add_feature_slide(slide6, "Intermediate Stoppages & Timetable Planner", "Feature 05", stop_info, f5_img if os.path.exists(f5_img) else mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 7: Feature 06 - Ticket Store & Departure Alarming Pin
    # -------------------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    feats6 = [
        ("🎟️", "Ticket Details Pinned Storage", "Store paper/online ticket photo, reference pin, train name, coach & seat info locally in app wallet."),
        ("🔔", "Journey Day Alarming Pin", "Triggers sound alarms & departure notifications prior to departure time at Kamalapur/CTG."),
        ("📲", "Offline Quick Pass Access", "Instant offline ticket reference retrieval at station entry gates without requiring mobile data.")
    ]
    add_feature_slide(slide7, "Ticket Store & Departure Alarming Pin", "Feature 06", feats6, mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 8: Major System Updates (100% Native PowerPoint Cards)
    # -------------------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8, LIGHT_BG)
    add_header(slide8, "Major System Updates — Multimodal Transit Hub", "Roadmap: Major")

    major_updates = [
        ("Dhaka Metro Rail (MRT-6 & Extension)", "Full integration with Dhaka Metro Rail Line 6 timetables, train interval headways, Rapid Pass balance check, and seamless station transfer guide at Kamalapur Intercity Junction.", "icon_metro_rail.png"),
        ("BusRouter BD (Inter-District Bus Schedules)", "Comprehensive highway bus network integration across Bangladesh (Hanif, Shyamoli, Shohag, Ena, Green Line) with counter locations, routes, and feeder bus connections to railway stations.", "icon_intercity_bus.png"),
        ("Launch & Waterway Router (Sadarghat & Riverine)", "Sadarghat launch terminal riverine schedule finder, deck/cabin fare lookup, and waterway routes for Southern Bangladesh districts (Barishal, Bhola, Chandpur, Patuakhali).", "icon_launch_router.png"),
        ("Unified Multi-Modal Trip Planner", "Door-to-door trip planner allowing commuters to calculate combined journey routes (Local Metro ➔ Intercity Train ➔ Highway Bus / River Launch) in a single unified interface.", "icon_multimodal_planner.png")
    ]

    for i, (m_title, m_desc, m_icon) in enumerate(major_updates):
        row = i // 2
        col = i % 2
        x = Inches(0.8) + col * Inches(5.95)
        y = Inches(1.55) + row * Inches(2.6)
        
        add_card(slide8, x, y, Inches(5.75), Inches(2.45), CARD_BG, border_color=CYAN)
        
        hdr_bar = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.75), Inches(0.48))
        hdr_bar.fill.solid()
        hdr_bar.fill.fore_color.rgb = NAVY
        hdr_bar.line.fill.background()
        tf_h = hdr_bar.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = m_title
        p_h.font.size = Pt(13.5)
        p_h.font.bold = True
        p_h.font.color.rgb = RGBColor(255, 255, 255)
        p_h.alignment = PP_ALIGN.CENTER

        icon_path = os.path.join(ASSETS_DIR, m_icon)
        if os.path.exists(icon_path):
            slide8.shapes.add_picture(icon_path, x + Inches(0.18), y + Inches(0.6), Inches(1.5), Inches(1.5))
            tb_left = x + Inches(1.8)
            tb_width = Inches(3.8)
        else:
            tb_left = x + Inches(0.2)
            tb_width = Inches(5.35)

        tb_card = slide8.shapes.add_textbox(tb_left, y + Inches(0.55), tb_width, Inches(1.8))
        tf_c = tb_card.text_frame
        tf_c.word_wrap = True
        p = tf_c.paragraphs[0]
        p.text = m_desc
        p.font.size = Pt(12.5)
        p.font.color.rgb = TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 9: Minor System Updates & Smart Enhancements (2x3 Grid)
    # -------------------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_bg(slide9, LIGHT_BG)
    add_header(slide9, "Minor Updates & Smart Commuter Enhancements", "Roadmap: Minor")

    minor_updates = [
        ("Live Train GPS Tracker", "Real-time GPS tracking on Google Maps showing live train position, current speed, and automated station arrival alerts.", "icon_gps_tracker.png"),
        ("Seat Facilities & Coach Visualizer", "Virtual coach layout showing seat power outlets, AC vent position, reclining angle, window view, and hygiene ratings.", "icon_seat_visualizer.png"),
        ("AI Schedule Delay Predictor", "Machine learning algorithm analyzing junction bottlenecks (Tongi, Joydebpur) to predict realistic delay arrival times.", "icon_ai_delay.png"),
        ("Station Offline Layout Maps", "Offline station map guides showing platforms 1-10, ticket counters, food court, parking, and luggage locker points.", "icon_station_map.png"),
        ("Coach Crowd Density Monitor", "Live passenger crowding indicators for each coach (Low, Moderate, High) to help passengers find spacious seating.", "icon_crowd_monitor.png"),
        ("Smart Commuter Hub", "Push notifications, fare change alerts & community delay reports for active passengers.", "icon_smart_hub.png")
    ]

    for i, (min_title, min_desc, min_icon) in enumerate(minor_updates):
        row = i // 3
        col = i % 3
        x = Inches(0.8) + col * Inches(3.98)
        y = Inches(1.55) + row * Inches(2.6)
        
        add_card(slide9, x, y, Inches(3.78), Inches(2.45), CARD_BG, border_color=CYAN)
        
        hdr_pill = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), y + Inches(0.15), Inches(3.48), Inches(0.38))
        hdr_pill.fill.solid()
        hdr_pill.fill.fore_color.rgb = NAVY
        hdr_pill.line.fill.background()
        tf_hp = hdr_pill.text_frame
        p_hp = tf_hp.paragraphs[0]
        p_hp.text = min_title
        p_hp.font.size = Pt(12.5)
        p_hp.font.bold = True
        p_hp.font.color.rgb = RGBColor(255, 255, 255)
        p_hp.alignment = PP_ALIGN.CENTER

        icon_path = os.path.join(ASSETS_DIR, min_icon)
        if os.path.exists(icon_path):
            slide9.shapes.add_picture(icon_path, x + Inches(0.15), y + Inches(0.6), Inches(1.2), Inches(1.2))
            tb_left = x + Inches(1.4)
            tb_width = Inches(2.25)
        else:
            tb_left = x + Inches(0.2)
            tb_width = Inches(3.38)

        tb_card = slide9.shapes.add_textbox(tb_left, y + Inches(0.58), tb_width, Inches(1.75))
        tf_c = tb_card.text_frame
        tf_c.word_wrap = True
        p = tf_c.paragraphs[0]
        p.text = min_desc
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 10: Conclusion & Q&A Slide (Dark Theme)
    # -------------------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_bg(slide10, LIGHT_BG)
    add_header(slide10, "Conclusion & Q&A", "Final Remarks")

    # Left Impact Card
    add_card(slide10, Inches(0.8), Inches(1.55), Inches(5.6), Inches(5.45), CARD_BG, border_color=NAVY)
    tb = slide10.shapes.add_textbox(Inches(1.0), Inches(1.75), Inches(5.2), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Bangladesh Railway Impact"
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = NAVY

    impacts = [
        ("🎯", "100% Transparent From-To Availability", "Delivers 100% transparent From-To train availability, seat counts, and seat class rates."),
        ("⚡", "Instant Seat-wise Fare Calculator", "Eliminates counter visit hassles by providing an instant seat-wise fare calculator on web/mobile."),
        ("🛡️", "Smart Alerts & Ticket Management", "Pinned ticket storage & departure alarm alerts ensure commuters never miss their train.")
    ]

    for icon_sym, title, desc in impacts:
        p = tf.add_paragraph()
        p.text = f"{icon_sym}  {title}"
        p.font.size = Pt(15.5)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_before = Pt(12)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13.5)
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(2)

    # Right Dark Q&A Box (Clean Redesigned Thank You Card without Badges)
    add_card(slide10, Inches(6.6), Inches(1.55), Inches(5.9), Inches(5.45), DARK_BG, border_color=CYAN)
    
    # Glowing Cyan Header Pill Badge (Clean Pure Text)
    ty_badge = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(2.1), Inches(4.7), Inches(1.0))
    ty_badge.fill.solid()
    ty_badge.fill.fore_color.rgb = CYAN
    ty_badge.line.fill.background()
    tf_ty = ty_badge.text_frame
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "THANK YOU!"
    p_ty.font.size = Pt(34)
    p_ty.font.bold = True
    p_ty.font.color.rgb = DARK_BG
    p_ty.alignment = PP_ALIGN.CENTER

    # Center Text Box for Q&A
    tb_r = slide10.shapes.add_textbox(Inches(7.0), Inches(3.4), Inches(5.1), Inches(3.3))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    
    p_r2 = tf_r.paragraphs[0]
    p_r2.text = "Questions & Answers Session"
    p_r2.font.size = Pt(22)
    p_r2.font.bold = True
    p_r2.font.color.rgb = RGBColor(255, 255, 255)
    p_r2.alignment = PP_ALIGN.CENTER

    p_r3 = tf_r.add_paragraph()
    p_r3.text = "We welcome any questions, feedback, and suggestions regarding TrainRouter BD System."
    p_r3.font.size = Pt(14.5)
    p_r3.font.color.rgb = RGBColor(203, 213, 225)
    p_r3.space_before = Pt(14)
    p_r3.alignment = PP_ALIGN.CENTER

    # Accent Divider Line
    div_line = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(5.3), Inches(2.5), Inches(0.04))
    div_line.fill.solid()
    div_line.fill.fore_color.rgb = CYAN
    div_line.line.fill.background()

    p_r4 = tf_r.add_paragraph()
    p_r4.text = "Bangladesh Railway Smart Journey Platform"
    p_r4.font.size = Pt(13.5)
    p_r4.font.bold = True
    p_r4.font.color.rgb = CYAN
    p_r4.space_before = Pt(32)
    p_r4.alignment = PP_ALIGN.CENTER

    target_files = ["TrainRouter_Features_Showcase.pptx", "TrainRouter_Features_Only.pptx", "TrainRouter_Features_Slide8_to_17.pptx"]
    for filename in target_files:
        try:
            prs.save(filename)
            print(f"Features presentation saved to: {os.path.abspath(filename)}")
        except Exception as e:
            print(f"Error saving {filename}: {e}")

if __name__ == "__main__":
    create_features_presentation()
