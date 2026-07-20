import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

def create_presentation():
    prs = Presentation()
    
    script_dir = os.path.dirname(os.path.abspath(__file__))
    ASSETS_DIR = os.path.join(script_dir, "assets")
    
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # blank slide layout

    # Color Palette
    NAVY = RGBColor(15, 76, 129)          # #0F4C81 - Primary Deep
    CYAN = RGBColor(0, 180, 216)          # #00B4D8 - Accent Teal/Cyan
    DARK_BG = RGBColor(11, 25, 44)        # #0B192C - Premium Dark BG
    LIGHT_BG = RGBColor(248, 250, 252)    # #F8FAFC - Clean Light BG
    CARD_BG = RGBColor(255, 255, 255)     # #FFFFFF
    TEXT_MAIN = RGBColor(30, 41, 59)      # #1E293B
    TEXT_MUTED = RGBColor(100, 116, 139)  # #64748B
    ACCENT_BLUE = RGBColor(224, 242, 254) # #E0F2FE

    def add_bg(slide, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_header(slide, title_text, category_text="TrainRouter BD", dark_theme=False):
        # Header bar badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.45), Inches(2.8), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = CYAN if dark_theme else NAVY
        badge.line.fill.background()
        tf = badge.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"•  {category_text.upper()}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = DARK_BG if dark_theme else RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Title
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(9.5), Inches(0.8))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255) if dark_theme else NAVY

        # Top Right Brand Badge
        brand_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.7), Inches(0.45), Inches(1.8), Inches(0.35))
        brand_badge.fill.solid()
        brand_badge.fill.fore_color.rgb = RGBColor(30, 41, 59) if dark_theme else RGBColor(241, 245, 249)
        brand_badge.line.color.rgb = CYAN if dark_theme else NAVY
        brand_badge.line.width = Pt(1)
        tf_b = brand_badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = "TrainRouter BD"
        p_b.font.size = Pt(11)
        p_b.font.bold = True
        p_b.font.color.rgb = CYAN if dark_theme else NAVY
        p_b.alignment = PP_ALIGN.CENTER

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # -------------------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1, DARK_BG)

    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(7.5), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "TRAINROUTER BD"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    p2 = tf.add_paragraph()
    p2.text = "BD Train Route Finder, Seat Availability & Fare Calculator"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.space_before = Pt(14)

    p3 = tf.add_paragraph()
    p3.text = "An intelligent software system engineered for Bangladesh Railway commuters to search trains by origin & destination, inspect seat classes (Shovon, Snigdha, AC Berth, AC Cabin), check live available seat counts, and calculate journey fares."
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_before = Pt(16)

    hero_path = os.path.join(ASSETS_DIR, "hero_train.png")
    if os.path.exists(hero_path):
        slide1.shapes.add_picture(hero_path, Inches(8.8), Inches(1.5), Inches(3.8), Inches(4.5))

    add_card(slide1, Inches(1.0), Inches(6.0), Inches(11.333), Inches(0.8), bg_color=NAVY)
    tb = slide1.shapes.add_textbox(Inches(1.2), Inches(6.15), Inches(11.0), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "TrainRouter BD Project Presentation | CSE Software Engineering"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # -------------------------------------------------------------------------
    # SLIDE 2: Team Members Slide
    # -------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2, LIGHT_BG)
    add_header(slide2, "Project Engineering Team", "Project Contributors")

    members = [
        {"img": os.path.join(ASSETS_DIR, "Hasib.jpeg"), "id": "ID: 00724205101098"},
        {"img": os.path.join(ASSETS_DIR, "Diganta.jpeg"), "id": "ID: 00724205101103"},
        {"img": os.path.join(ASSETS_DIR, "imran.jpeg"), "id": "ID: 00724205101123"}
    ]

    card_w = Inches(3.5)
    card_h = Inches(4.5)
    start_x = Inches(0.9)
    gap = Inches(0.6)

    for i, m in enumerate(members):
        x = start_x + i * (card_w + gap)
        y = Inches(1.8)
        
        card = add_card(slide2, x, y, card_w, card_h, CARD_BG, border_color=CYAN)
        
        if os.path.exists(m["img"]):
            slide2.shapes.add_picture(m["img"], x + Inches(0.55), y + Inches(0.5), Inches(2.4), Inches(2.4))
        
        badge_shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), y + Inches(3.2), Inches(2.9), Inches(0.7))
        badge_shape.fill.solid()
        badge_shape.fill.fore_color.rgb = ACCENT_BLUE
        badge_shape.line.color.rgb = NAVY
        badge_shape.line.width = Pt(1.5)
        
        tf = badge_shape.text_frame
        p = tf.paragraphs[0]
        p.text = m["id"]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------------------
    # SLIDE 3: System Vision / What is TrainRouter BD?
    # -------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3, LIGHT_BG)
    add_header(slide3, "What is TrainRouter BD?", "System Vision")

    points = [
        ("From-To Station Search Engine", "Enables Bangladesh Railway commuters to search available intercity trains between any origin and destination (e.g. Dhaka to Cox's Bazar, Sylhet, CTG)."),
        ("Live Seat Inventory & Category Lookup", "Displays real-time available seat counts for all seat classes: Shovon (S_CHAIR), Snigdha (AC Chair), AC Berth (AC_BATH), and AC Cabin (AC_CABIN)."),
        ("Seat-Wise Fare Calculator Engine", "Computes total journey cost dynamically in BDT (৳) based on selected seat class, distance tariff matrix, and passenger quantity (No ticket purchasing).")
    ]

    for i, (title, desc) in enumerate(points):
        y = Inches(1.8) + i * Inches(1.7)
        add_card(slide3, Inches(0.9), y, Inches(6.5), Inches(1.45), CARD_BG, border_color=NAVY)
        
        tb = slide3.shapes.add_textbox(Inches(1.1), y + Inches(0.15), Inches(6.1), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"0{i+1}. {title}"
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = NAVY
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(4)

    map_img = os.path.join(ASSETS_DIR, "train_route_map.png")
    if os.path.exists(map_img):
        slide3.shapes.add_picture(map_img, Inches(7.7), Inches(1.8), Inches(4.7), Inches(4.8))

    # -------------------------------------------------------------------------
    # SLIDE 4: Problem Statement
    # -------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4, LIGHT_BG)
    add_header(slide4, "Challenges in Railway Commute Planning", "Background Context")

    problems = [
        ("Uncertain Train Schedules", "Passengers often lack instant visibility on how many intercity trains run between stations and their departure/arrival schedules."),
        ("Lack of Seat Class Details", "Commuters struggle to inspect real-time available seat counts for specific classes (Shovon, Snigdha, AC Berth, AC Cabin) before traveling."),
        ("Complex Fare Breakdown", "Calculating total journey cost for family or group travel across different seat categories requires tedious manual computations.")
    ]

    for i, (p_title, p_desc) in enumerate(problems):
        x = Inches(0.9) + i * Inches(3.9)
        add_card(slide4, x, Inches(1.8), Inches(3.6), Inches(4.8), CARD_BG, border_color=RGBColor(239, 68, 68))
        
        tb = slide4.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.2), Inches(4.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(185, 28, 28)
        
        p2 = tf.add_paragraph()
        p2.text = p_desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(14)

    # -------------------------------------------------------------------------
    # SLIDE 5: Proposed Solution & Flow
    # -------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5, LIGHT_BG)
    add_header(slide5, "How TrainRouter BD Works", "System Architecture")

    pillars = [
        ("Step 1: From-To Input", "Select origin and destination BD stations with preferred journey date."),
        ("Step 2: Database Query", "Queries Bangladesh Railway schedules for active intercity trains on the route."),
        ("Step 3: Seat Inventory Lookup", "Retrieves available seat counts for Shovon, Snigdha, AC Berth, and AC Cabin."),
        ("Step 4: Seat-Wise Fare Calculation", "Computes fare per seat class in BDT (৳) and total journey cost for N passengers.")
    ]

    for i, (title, desc) in enumerate(pillars):
        row = i // 2
        col = i % 2
        x = Inches(0.9) + col * Inches(5.9)
        y = Inches(1.8) + row * Inches(2.5)
        
        add_card(slide5, x, y, Inches(5.6), Inches(2.2), CARD_BG, border_color=CYAN)
        
        tb = slide5.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(5.0), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = NAVY
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(8)

    # -------------------------------------------------------------------------
    # SLIDE 6: Bangladesh Railway Data Schema
    # -------------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6, LIGHT_BG)
    add_header(slide6, "Bangladesh Railway Relational Schema", "Data Model")

    entities = [
        ("BD Stations Catalog", "station_id, station_name (Dhaka Kamalapur, CTG, Sylhet, CXB), station_code, zone"),
        ("Intercity Schedules", "train_id, train_name (Suborno, Parabat, Sonar Bangla), origin, destination, off_day"),
        ("Seat Inventory Pool", "train_id, seat_type (Shovon, Snigdha, AC_Berth, AC_Cabin), available_seats_count"),
        ("BR Fare Tariff Rules", "route_km, base_tariff_rate_bdt, seat_class_multiplier, govt_vat_percentage (15%)")
    ]

    for i, (e_title, e_fields) in enumerate(entities):
        y = Inches(1.8) + i * Inches(1.25)
        add_card(slide6, Inches(0.9), y, Inches(11.533), Inches(1.05), CARD_BG, border_color=NAVY)
        
        tb = slide6.shapes.add_textbox(Inches(1.2), y + Inches(0.12), Inches(11.0), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{e_title}: "
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        
        p_run = p.add_run()
        p_run.text = e_fields
        p_run.font.size = Pt(14)
        p_run.font.bold = False
        p_run.font.color.rgb = TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 7: BD Dynamic Fare Computation
    # -------------------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7, LIGHT_BG)
    add_header(slide7, "Bangladesh Railway Fare & Multiplier Engine", "Calculation Rules")

    add_card(slide7, Inches(0.9), Inches(1.8), Inches(5.6), Inches(4.8), CARD_BG, border_color=NAVY)
    tb = slide7.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.2), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "BD Fare Formula (BDT ৳)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    p2 = tf.add_paragraph()
    p2.text = "• Seat Fare = (Base Tariff Rate × Route Distance KM) × Class Multiplier + 15% Govt VAT.\n• Total Journey Fare = Single Seat Fare × Passenger Count.\n• Distance matrices computed across Eastern & Western Rail Tracks."
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MAIN
    p2.space_before = Pt(10)

    add_card(slide7, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), CARD_BG, border_color=CYAN)
    tb2 = slide7.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.3))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "BR Seat Class Multipliers"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    p2 = tf2.add_paragraph()
    p2.text = "• Shovon Chair (S_CHAIR): 1.0x (Standard Non-AC Padded Recliner Base Tariff)\n• Snigdha (AC Chair): 1.8x (Air-Conditioned Premium Recliner Class)\n• AC Berth (AC_BATH): 2.5x (Night Sleeper Berth Class)\n• AC Cabin (AC_CABIN): 2.8x (Private Compartment Cabin)"
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MAIN
    p2.space_before = Pt(10)

    BADGE_COLORS = [
        RGBColor(15, 76, 129),    # Deep Navy
        RGBColor(0, 150, 199),   # Vibrant Cyan/Teal
        RGBColor(16, 124, 65),   # Emerald Green
        RGBColor(217, 119, 6),   # Amber Gold
        RGBColor(124, 58, 237),  # Royal Purple
    ]

    def add_feature_slide(slide, title_text, category_text, points_list, img_path):
        add_bg(slide, LIGHT_BG)
        add_header(slide, title_text, category_text)

        # Left Content Container Card
        add_card(slide, Inches(0.8), Inches(1.55), Inches(5.6), Inches(5.45), CARD_BG, border_color=NAVY)
        
        tb = slide.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(5.3), Inches(5.15))
        tf = tb.text_frame
        tf.word_wrap = True
        
        first = True
        for idx, item in enumerate(points_list):
            if len(item) == 3:
                icon_sym, t_text, d_text = item
            elif len(item) == 2:
                icon_sym = "►"
                t_text, d_text = item
            else:
                icon_sym = "►"
                t_text, d_text = f"Feature Point 0{idx+1}", item[0]

            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            
            p.text = f"{icon_sym}  {t_text}"
            p.font.size = Pt(15.5)
            p.font.bold = True
            p.font.color.rgb = BADGE_COLORS[idx % len(BADGE_COLORS)]
            p.space_before = Pt(6) if idx > 0 else Pt(0)
            
            p_desc = tf.add_paragraph()
            p_desc.text = d_text
            p_desc.font.size = Pt(13)
            p_desc.font.color.rgb = TEXT_MAIN
            p_desc.space_after = Pt(6)

        # Right Image Frame Card (Expanded Wide Frame)
        card_left = Inches(6.6)
        card_top = Inches(1.55)
        card_width = Inches(5.9)
        card_height = Inches(5.45)
        add_card(slide, card_left, card_top, card_width, card_height, RGBColor(241, 245, 249), border_color=CYAN)
        
        if img_path and os.path.exists(img_path):
            with Image.open(img_path) as img:
                orig_w, orig_h = img.size
            aspect = orig_w / float(orig_h)
            
            # Calculate optimal wide dimensions preserving natural aspect ratio
            max_img_w = Inches(5.6)
            max_img_h = Inches(5.1)
            
            calc_h = max_img_w / aspect
            if calc_h <= max_img_h:
                final_w = max_img_w
                final_h = calc_h
            else:
                final_h = max_img_h
                final_w = final_h * aspect
                
            img_left = card_left + (card_width - final_w) / 2.0
            img_top = card_top + (card_height - final_h) / 2.0
            
            slide.shapes.add_picture(img_path, img_left, img_top, final_w, final_h)

    mockup_img_v2 = os.path.join(ASSETS_DIR, "train_app_mockup_v2.jpg")
    mockup_img = mockup_img_v2 if os.path.exists(mockup_img_v2) else os.path.join(ASSETS_DIR, "train_app_mockup.png")

    # -------------------------------------------------------------------------
    # SLIDE 8: System Features Showcase Overview
    # -------------------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    feats_overview = [
        ("🔍", "From-To Station Train Search", "Select origin, destination & date to view all active running trains."),
        ("📊", "Available Seat Counts per Class", "Live inventory breakdown for Shovon, Snigdha, AC Berth, and AC Cabin."),
        ("🧮", "Seat-Wise Fare Calculator (BDT ৳)", "Instant cost lookup for single seat & multi-passenger journeys."),
        ("📍", "Intermediate Stoppages & Timetables", "Arrival/departure times at major junction stations."),
        ("📌", "Pinned Ticket Store & Departure Alarm", "Store offline ticket reference & set departure alert notifications.")
    ]
    add_feature_slide(slide8, "System Features Showcase Overview", "Features Overview", feats_overview, mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 9: Feature 01 - From-To & Date-Based Search
    # -------------------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    feats1 = [
        ("🔎", "Smart Autocomplete Station Lookup", "Select From (Dhaka Kamalapur) and To (Chattogram, Sylhet, Cox's Bazar, Rajshahi)."),
        ("📅", "Journey Date Picker & Off-Day Filter", "Pick any travel date. System automatically checks weekly off-days and filters non-operating trains."),
        ("🏷️", "Available Train Counter Badge", "Instantly shows top header summary: '4 Intercity Trains Available on Selected Date'.")
    ]
    add_feature_slide(slide9, "From-To & Date-Based Train Search", "Feature 01", feats1, mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 10: Feature 02 - Date-Wise Available Train List & Seat Count
    # -------------------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    feats2 = [
        ("🚆", "Date-Wise Active Train Schedule", "Displays exact running intercity trains on selected date (Suborno 701, Parabat 709, Cox's Bazar Express 813)."),
        ("💺", "Real-Time Seat Inventory per Class", "Live seat count breakdown (Shovon: 120 seats, Snigdha: 45 seats, AC Berth: 12 seats)."),
        ("⏰", "Train Timetable Details", "Exact departure time, destination arrival time, and estimated total journey duration.")
    ]
    add_feature_slide(slide10, "Date-Wise Active Train Schedule & Seats", "Feature 02", feats2, mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 11: Feature 03 - Seat Categories & Specifications
    # -------------------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    seats_info = [
        ("💺", "Shovon Chair (S_CHAIR)", "Standard non-AC padded recliner seat with 1.0x base tariff rate."),
        ("❄️", "Snigdha (AC Chair)", "Premium air-conditioned recliner seating with 1.8x base tariff rate."),
        ("🌙", "AC Berth (AC_BATH)", "Night sleeper berth compartment for long-distance overnight travel (2.5x tariff)."),
        ("👑", "AC Cabin (AC_CABIN)", "Private lockable compartment cabin suitable for family travel (2.8x tariff).")
    ]
    f3_img = os.path.join(ASSETS_DIR, "feature_03_seat_types.png")
    add_feature_slide(slide11, "Available Seat Types & Categories", "Feature 03", seats_info, f3_img if os.path.exists(f3_img) else mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 12: Feature 04 - Seat-Wise Fare Calculator
    # -------------------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    fare_info = [
        ("📏", "Dynamic Distance Fare Calculation", "Automatically looks up route distance (KM) between selected stations."),
        ("🏷️", "Class Multipliers & VAT Addition", "Applies exact seat multiplier (1.0x to 2.8x) + 15% Bangladesh Govt VAT."),
        ("👨‍👩‍👧‍👦", "Group / Family Trip Costing", "Calculates total cost for N passengers dynamically without requiring ticket booking.")
    ]
    fare_user_img = os.path.join(ASSETS_DIR, "fare_calc_user_mockup.jpg")
    add_feature_slide(slide12, "Seat-Wise Fare Calculator (BDT ৳)", "Feature 04", fare_info, fare_user_img if os.path.exists(fare_user_img) else mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 13: Feature 05 - Stoppage & Timetable Planner
    # -------------------------------------------------------------------------
    slide13 = prs.slides.add_slide(blank_layout)
    stop_info = [
        ("🚉", "Stoppage Junction Timetables", "Detailed arrival & departure times at intermediate stations (Airport, Joydebpur, Akhaura)."),
        ("🌐", "Cross-Zone Rail Connectivity", "Route mapping connecting Dhaka, Rajshahi, Khulna, Sylhet, and Cox's Bazar."),
        ("🛤️", "Gauge Compatibility Indicator", "Clearly demarcates Meter Gauge (MG), Broad Gauge (BG), and Dual Gauge (DG) lines.")
    ]
    f5_img = os.path.join(ASSETS_DIR, "feature_05_route_stoppages.png")
    add_feature_slide(slide13, "Intermediate Stoppages & Timetable Planner", "Feature 05", stop_info, f5_img if os.path.exists(f5_img) else mockup_img)

    # -------------------------------------------------------------------------
    # SLIDE 14: Feature 06 - Ticket Store & Departure Alarming Pin
    # -------------------------------------------------------------------------
    slide14 = prs.slides.add_slide(blank_layout)
    feats6 = [
        ("🎟️", "Ticket Details Pinned Storage", "Store paper/online ticket photo, reference pin, train name, coach & seat info locally in app wallet."),
        ("🔔", "Journey Day Alarming Pin", "Triggers sound alarms & departure notifications prior to departure time at Kamalapur/CTG."),
        ("📲", "Offline Quick Pass Access", "Instant offline ticket reference retrieval at station entry gates without requiring mobile data.")
    ]
    add_feature_slide(slide14, "Ticket Store & Departure Alarming Pin", "Feature 06", feats6, mockup_img)

    if os.path.exists(mockup_img):
        slide14.shapes.add_picture(mockup_img, Inches(7.4), Inches(1.8), Inches(5.0), Inches(5.0))

    # -------------------------------------------------------------------------
    # SLIDE 15: Major System Updates (100% Native PowerPoint Cards)
    # -------------------------------------------------------------------------
    slide15 = prs.slides.add_slide(blank_layout)
    add_bg(slide15, LIGHT_BG)
    add_header(slide15, "Major System Updates — Multimodal Transit Hub", "Roadmap: Major")

    major_updates = [
        ("Dhaka Metro Rail (MRT-6 & Extension)", "Full integration with Dhaka Metro Rail Line 6 timetables, train interval headways, Rapid Pass balance check, and seamless station transfer guide at Kamalapur Intercity Junction.", "icon_metro_rail.png"),
        ("BusRouter BD (Inter-District Bus Schedules)", "Comprehensive highway bus network integration across Bangladesh (Hanif, Shyamoli, Shohag, Ena, Green Line) with counter locations, routes, and feeder bus connections to railway stations.", "icon_intercity_bus.png"),
        ("Launch & Waterway Router (Sadarghat & Riverine)", "Sadarghat launch terminal riverine schedule finder, deck/cabin fare lookup, and waterway routes for Southern Bangladesh districts (Barishal, Bhola, Chandpur, Patuakhali).", "icon_launch_router.png"),
        ("Unified Multi-Modal Trip Planner", "Door-to-door trip planner allowing commuters to calculate combined journey routes (Local Metro ➔ Intercity Train ➔ Highway Bus / River Launch) in a single unified interface.", "icon_multimodal_planner.png")
    ]

    for i, (m_title, m_desc, m_icon) in enumerate(major_updates):
        row = i // 2
        col = i % 2
        x = Inches(0.9) + col * Inches(5.9)
        y = Inches(1.8) + row * Inches(2.5)
        
        add_card(slide15, x, y, Inches(5.6), Inches(2.3), CARD_BG, border_color=CYAN)
        
        hdr_bar = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.6), Inches(0.45))
        hdr_bar.fill.solid()
        hdr_bar.fill.fore_color.rgb = NAVY
        hdr_bar.line.fill.background()
        tf_h = hdr_bar.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = f"  {m_title}"
        p_h.font.size = Pt(14)
        p_h.font.bold = True
        p_h.font.color.rgb = RGBColor(255, 255, 255)
        
        icon_path = os.path.join(ASSETS_DIR, m_icon)
        if os.path.exists(icon_path):
            slide15.shapes.add_picture(icon_path, x + Inches(0.18), y + Inches(0.55), Inches(1.4), Inches(1.4))
            tb_left = x + Inches(1.68)
            tb_width = Inches(3.72)
        else:
            tb_left = x + Inches(0.2)
            tb_width = Inches(5.2)

        tb = slide15.shapes.add_textbox(tb_left, y + Inches(0.5), tb_width, Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = m_desc
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 16: Minor System Updates & Smart Enhancements (100% Native PPT Cards)
    # -------------------------------------------------------------------------
    slide16 = prs.slides.add_slide(blank_layout)
    add_bg(slide16, LIGHT_BG)
    add_header(slide16, "Minor Updates & Smart Commuter Enhancements", "Roadmap: Minor")

    minor_updates = [
        ("Live Train GPS Tracker", "Real-time GPS tracking on Google Maps showing live train position, current speed, and automated station arrival alerts.", "icon_gps_tracker.png"),
        ("Seat Facilities & Coach Visualizer", "Virtual coach layout showing seat power outlets, AC vent position, reclining angle, window view, and hygiene ratings.", "icon_seat_visualizer.png"),
        ("AI Schedule Delay Predictor", "Machine learning algorithm analyzing junction bottlenecks (Tongi, Joydebpur) to predict realistic delay arrival times.", "icon_ai_delay.png"),
        ("Station Offline Layout Maps", "Offline station map guides showing platforms 1-10, ticket counters, food court, parking, and luggage locker points.", "icon_station_map.png"),
        ("Coach Crowd Density Monitor", "Live passenger crowding indicators for each coach (Low, Moderate, High) to help passengers find spacious seating.", "icon_crowd_monitor.png"),
        ("Smart Commuter Hub", "Push notifications, fare change alerts & community delay reports for active passengers.", "icon_smart_hub.png")
    ]

    for i, (min_title, min_desc, min_icon) in enumerate(minor_updates):
        row = i // 3
        col = i % 3
        x = Inches(0.9) + col * Inches(3.95)
        y = Inches(1.8) + row * Inches(2.5)
        
        add_card(slide16, x, y, Inches(3.75), Inches(2.3), CARD_BG, border_color=CYAN)
        
        hdr_pill = slide16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), y + Inches(0.15), Inches(3.45), Inches(0.35))
        hdr_pill.fill.solid()
        hdr_pill.fill.fore_color.rgb = NAVY
        hdr_pill.line.fill.background()
        tf_hp = hdr_pill.text_frame
        p_hp = tf_hp.paragraphs[0]
        p_hp.text = min_title
        p_hp.font.size = Pt(12)
        p_hp.font.bold = True
        p_hp.font.color.rgb = RGBColor(255, 255, 255)
        p_hp.alignment = PP_ALIGN.CENTER

        icon_path = os.path.join(ASSETS_DIR, min_icon)
        if os.path.exists(icon_path):
            slide16.shapes.add_picture(icon_path, x + Inches(0.15), y + Inches(0.58), Inches(1.1), Inches(1.1))
            tb_left = x + Inches(1.3)
            tb_width = Inches(2.3)
        else:
            tb_left = x + Inches(0.15)
            tb_width = Inches(3.45)

        tb = slide16.shapes.add_textbox(tb_left, y + Inches(0.55), tb_width, Inches(1.65))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = min_desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 17: Conclusion & Q&A (100% Native PowerPoint Layout)
    # -------------------------------------------------------------------------
    slide17 = prs.slides.add_slide(blank_layout)
    add_bg(slide17, LIGHT_BG)
    add_header(slide17, "Conclusion & Q&A", "Final Remarks")

    add_card(slide17, Inches(0.9), Inches(1.8), Inches(6.0), Inches(4.9), CARD_BG, border_color=NAVY)
    
    tb = slide17.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.6), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Bangladesh Railway Impact"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    impacts = [
        ("100% Transparent From-To Availability", "Delivers 100% transparent From-To train availability, seat counts, and seat class rates."),
        ("Instant Seat-wise Fare Calculator", "Eliminates counter visit hassles by providing an instant seat-wise fare calculator on web/mobile."),
        ("Smart Alerts & Ticket Management", "Pinned ticket storage & departure alarm alerts ensure commuters never miss their train.")
    ]
    
    for imp_title, imp_desc in impacts:
        p1 = tf.add_paragraph()
        p1.text = f"•  {imp_title}"
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        p1.space_before = Pt(12)
        
        p2 = tf.add_paragraph()
        p2.text = f"    {imp_desc}"
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_MAIN
        p2.space_before = Pt(4)

    # Right Dark Q&A Box (Clean Redesigned Thank You Card without Badges)
    add_card(slide17, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1), DARK_BG, border_color=CYAN)
    
    # Glowing Cyan Header Pill Badge (Clean Pure Text)
    ty_badge = slide17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(2.2), Inches(4.7), Inches(0.95))
    ty_badge.fill.solid()
    ty_badge.fill.fore_color.rgb = CYAN
    ty_badge.line.fill.background()
    tf_ty = ty_badge.text_frame
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "THANK YOU!"
    p_ty.font.size = Pt(32)
    p_ty.font.bold = True
    p_ty.font.color.rgb = DARK_BG
    p_ty.alignment = PP_ALIGN.CENTER

    # Center Text Box for Q&A
    tb_r = slide17.shapes.add_textbox(Inches(7.1), Inches(3.4), Inches(5.1), Inches(3.1))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    
    p_r2 = tf_r.paragraphs[0]
    p_r2.text = "Questions & Answers Session"
    p_r2.font.size = Pt(22)
    p_r2.font.bold = True
    p_r2.font.color.rgb = RGBColor(255, 255, 255)
    p_r2.alignment = PP_ALIGN.CENTER

    p_r3 = tf_r.add_paragraph()
    p_r3.text = "We welcome any questions, feedback, and suggestions regarding TrainRouter BD System."
    p_r3.font.size = Pt(14)
    p_r3.font.color.rgb = RGBColor(203, 213, 225)
    p_r3.space_before = Pt(14)
    p_r3.alignment = PP_ALIGN.CENTER

    # Accent Divider Line
    div_line = slide17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.4), Inches(5.2), Inches(2.5), Inches(0.04))
    div_line.fill.solid()
    div_line.fill.fore_color.rgb = CYAN
    div_line.line.fill.background()

    p_r4 = tf_r.add_paragraph()
    p_r4.text = "Bangladesh Railway Smart Journey Platform"
    p_r4.font.size = Pt(13)
    p_r4.font.bold = True
    p_r4.font.color.rgb = CYAN
    p_r4.space_before = Pt(30)
    p_r4.alignment = PP_ALIGN.CENTER

    target_files = ["TrainRouter_Presentation_Final.pptx", "TrainRouter_Presentation_Fixed.pptx", "TrainRouter_Presentation_v2.pptx", "TrainRouter_Presentation_v3.pptx"]
    saved = False
    for filename in target_files:
        try:
            prs.save(filename)
            print(f"Presentation successfully created at: {os.path.abspath(filename)}")
            saved = True
            break
        except PermissionError:
            continue
    if not saved:
        import time
        alt = f"TrainRouter_Presentation_{int(time.time())}.pptx"
        prs.save(alt)
        print(f"Saved presentation to alternate path: {os.path.abspath(alt)}")

if __name__ == "__main__":
    create_presentation()
